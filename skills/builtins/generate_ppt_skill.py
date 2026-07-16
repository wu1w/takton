"""
PPT 生成 Skill
根据用户提供的主题和大纲，生成 PowerPoint 文件
"""

import json
import logging
import os
import uuid
from typing import Any

from backend.skills.base import BaseSkill

logger = logging.getLogger(__name__)

PPT_OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "uploads", "ppt"))

SYSTEM_PROMPT = """你是一位专业的PPT制作专家。你的任务是根据用户提供的主题和要求，生成一份结构清晰、内容专业的PPT大纲和内容。

请按以下JSON格式输出PPT结构：
{
  "title": "PPT标题",
  "subtitle": "副标题",
  "slides": [
    {
      "title": "幻灯片标题",
      "content": ["要点1", "要点2", "要点3"],
      "layout": "title_content"
    }
  ]
}

布局类型(layout)：
- title_slide: 标题页
- title_content: 标题+内容
- two_content: 两栏内容
- section_header: 章节标题页
- comparison: 对比页

要求：
1. 内容专业、简洁，适合演示
2. 每页要点不超过5条
3. 总页数控制在10-20页
4. 输出必须是合法的JSON格式
"""


class GeneratePPTSkill(BaseSkill):
    """生成PPT文件"""

    name = "generate_ppt"
    description = "根据主题生成专业PPT演示文稿。需要主题、页数要求等参数。"
    parameters = {
        "type": "object",
        "properties": {
            "topic": {
                "type": "string",
                "description": "PPT主题",
            },
            "pages": {
                "type": "integer",
                "description": "期望页数（默认10页）",
                "default": 10,
            },
            "audience": {
                "type": "string",
                "description": "目标受众，如'技术团队'、'管理层'、'客户'",
                "default": "通用受众",
            },
            "outline": {
                "type": "string",
                "description": "可选的大纲或要点，用换行分隔",
                "default": "",
            },
        },
        "required": ["topic"],
    }

    async def execute(self, topic: str, pages: int = 10, audience: str = "通用受众", outline: str = "") -> str:
        """生成PPT文件"""
        try:
            from backend.services.llm import LLMServiceFactory

            # 构建提示词
            prompt = f"请为以下主题生成PPT内容：\n\n主题：{topic}\n目标受众：{audience}\n期望页数：{pages}页\n"
            if outline:
                prompt += f"大纲要点：\n{outline}\n\n"
            prompt += f"\n请严格按照以下系统提示生成PPT结构：\n{SYSTEM_PROMPT}"

            llm_service = LLMServiceFactory.get_service()
            response = ""
            async for chunk in llm_service.chat([{"role": "user", "content": prompt}], stream=False):
                response += chunk.delta or ""
                if chunk.finish_reason:
                    break

            # 提取JSON
            ppt_data = self._extract_json(response)
            if not ppt_data:
                return f"[Error] 无法解析PPT结构。LLM响应：{response[:500]}"

            # 生成PPT文件
            file_path = self._generate_ppt_file(ppt_data)
            file_name = os.path.basename(file_path)

            return (
                f"[Success] PPT生成完成！\n"
                f"标题：{ppt_data.get('title', topic)}\n"
                f"页数：{len(ppt_data.get('slides', []))}页\n"
                f"下载链接：/uploads/ppt/{file_name}\n"
                f"文件路径：{file_path}"
            )

        except Exception as e:
            logger.error(f"PPT generation failed: {e}")
            return f"[Error] PPT生成失败: {e}"

    def _extract_json(self, text: str) -> dict[str, Any] | None:
        """从LLM响应中提取JSON"""
        # 尝试找 ```json 代码块
        import re
        code_block = re.search(r"```(?:json)?\s*(.*?)\s*```", text, re.DOTALL)
        if code_block:
            text = code_block.group(1)

        # 尝试找最外层的大括号
        start = text.find("{")
        end = text.rfind("}")
        if start != -1 and end != -1 and end > start:
            try:
                return json.loads(text[start:end+1])
            except json.JSONDecodeError:
                pass
        return None

    def _generate_ppt_file(self, ppt_data: dict[str, Any]) -> str:
        """使用python-pptx生成PPT文件"""
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            # 如果没有python-pptx，生成markdown文件作为fallback
            return self._generate_markdown_fallback(ppt_data)

        prs = Presentation()
        title = ppt_data.get("title", "Untitled")
        subtitle = ppt_data.get("subtitle", "")
        slides = ppt_data.get("slides", [])

        # 如果没有slides，创建一个默认的
        if not slides:
            slides = [{"title": title, "content": ["内容待补充"], "layout": "title_content"}]

        for i, slide_data in enumerate(slides):
            layout_name = slide_data.get("layout", "title_content")
            slide_title = slide_data.get("title", "")
            content = slide_data.get("content", [])

            # 选择布局
            if layout_name == "title_slide" or i == 0:
                slide_layout = prs.slide_layouts[0]  # Title Slide
            elif layout_name == "section_header":
                slide_layout = prs.slide_layouts[2]  # Section Header
            elif layout_name == "two_content":
                slide_layout = prs.slide_layouts[5]  # Two Content
            else:
                slide_layout = prs.slide_layouts[1]  # Title and Content

            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            if slide.shapes.title:
                slide.shapes.title.text = slide_title

            # 设置内容
            if layout_name == "title_slide" and i == 0:
                # 标题页设置副标题
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        shape.text = subtitle or ""
            elif layout_name == "two_content" and content and len(slide.placeholders) > 2:
                # 双栏布局：将内容平均拆分到左右两栏
                mid = (len(content) + 1) // 2
                left_items, right_items = content[:mid], content[mid:]
                for idx, items in ((1, left_items), (2, right_items)):
                    body_shape = slide.placeholders[idx]
                    tf = body_shape.text_frame
                    tf.clear()
                    for j, item in enumerate(items):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        p.text = f"• {item}"
                        p.level = 0
                        p.font.size = Pt(18)
            elif content and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for j, item in enumerate(content):
                    if j == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = f"• {item}"
                    p.level = 0
                    p.font.size = Pt(18)

        # 保存文件
        os.makedirs(PPT_OUTPUT_DIR, exist_ok=True)
        file_name = f"{uuid.uuid4().hex}.pptx"
        file_path = os.path.join(PPT_OUTPUT_DIR, file_name)
        prs.save(file_path)
        return file_path

    def _generate_markdown_fallback(self, ppt_data: dict[str, Any]) -> str:
        """没有python-pptx时的fallback：生成markdown"""
        title = ppt_data.get("title", "Untitled")
        subtitle = ppt_data.get("subtitle", "")
        slides = ppt_data.get("slides", [])

        lines = [f"# {title}", f"\n> {subtitle}\n"]
        for slide in slides:
            lines.append(f"\n## {slide.get('title', '')}")
            for item in slide.get("content", []):
                lines.append(f"- {item}")

        os.makedirs(PPT_OUTPUT_DIR, exist_ok=True)
        file_name = f"{uuid.uuid4().hex}.md"
        file_path = os.path.join(PPT_OUTPUT_DIR, file_name)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
        return file_path
